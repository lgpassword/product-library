import os
import io
import re
import uuid
import zipfile
from datetime import datetime

from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Query, BackgroundTasks
from fastapi.responses import FileResponse, Response
from fastapi.staticfiles import StaticFiles
from urllib.parse import quote
from pydantic import BaseModel
from typing import Optional, List

import openpyxl
from openpyxl.styles import Font
from openpyxl.drawing.image import Image as XLImage

from database import get_conn, init_db, BASE_DIR, IMAGE_DIR, FILE_DIR, LIBRARY_DIR, DB_PATH

init_db()

app = FastAPI(title="产品库系统")


def now():
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def row_to_dict(row):
    return dict(row) if row is not None else None


def get_product_tags(conn, product_id):
    rows = conn.execute(
        "SELECT t.id, t.name FROM tag t JOIN product_tag pt ON t.id=pt.tag_id "
        "WHERE pt.product_id=? ORDER BY t.name", (product_id,)
    ).fetchall()
    return [dict(r) for r in rows]


def product_full(conn, p):
    if p is None:
        return None
    d = dict(p)
    d["tags"] = get_product_tags(conn, d["id"])
    atts = conn.execute(
        "SELECT id, filename, kind, size, created_at FROM attachment "
        "WHERE product_id=? ORDER BY id", (d["id"],)
    ).fetchall()
    d["attachments"] = [dict(a) for a in atts]
    return d


# ---------- 字典 ----------

def list_dict(conn, table):
    return [dict(r) for r in conn.execute(f"SELECT * FROM {table} ORDER BY id").fetchall()]


@app.get("/api/categories")
def get_categories():
    conn = get_conn()
    data = list_dict(conn, "category")
    conn.close()
    return data


@app.post("/api/categories")
def add_category(payload: dict):
    name = (payload.get("name") or "").strip()
    if not name:
        raise HTTPException(400, "名称不能为空")
    conn = get_conn()
    try:
        cur = conn.execute("INSERT INTO category(name) VALUES(?)", (name,))
        conn.commit()
        return {"id": cur.lastrowid, "name": name}
    except Exception as e:
        raise HTTPException(400, "已存在或参数错误")
    finally:
        conn.close()


@app.delete("/api/categories/{cid}")
def del_category(cid: int):
    conn = get_conn()
    conn.execute("DELETE FROM category WHERE id=?", (cid,))
    conn.commit()
    conn.close()
    return {"ok": True}


@app.get("/api/companies")
def get_companies():
    conn = get_conn()
    data = list_dict(conn, "company")
    conn.close()
    return data


@app.post("/api/companies")
def add_company(payload: dict):
    name = (payload.get("name") or "").strip()
    if not name:
        raise HTTPException(400, "名称不能为空")
    conn = get_conn()
    try:
        cur = conn.execute("INSERT INTO company(name) VALUES(?)", (name,))
        conn.commit()
        return {"id": cur.lastrowid, "name": name}
    except Exception:
        raise HTTPException(400, "已存在或参数错误")
    finally:
        conn.close()


@app.delete("/api/companies/{cid}")
def del_company(cid: int):
    conn = get_conn()
    conn.execute("DELETE FROM company WHERE id=?", (cid,))
    conn.commit()
    conn.close()
    return {"ok": True}


@app.get("/api/tags")
def get_tags():
    conn = get_conn()
    data = list_dict(conn, "tag")
    conn.close()
    return data


@app.post("/api/tags")
def add_tag(payload: dict):
    name = (payload.get("name") or "").strip()
    if not name:
        raise HTTPException(400, "名称不能为空")
    conn = get_conn()
    try:
        cur = conn.execute("INSERT INTO tag(name) VALUES(?)", (name,))
        conn.commit()
        return {"id": cur.lastrowid, "name": name}
    except Exception:
        raise HTTPException(400, "已存在或参数错误")
    finally:
        conn.close()


@app.delete("/api/tags/{tid}")
def del_tag(tid: int):
    conn = get_conn()
    conn.execute("DELETE FROM tag WHERE id=?", (tid,))
    conn.execute("DELETE FROM product_tag WHERE tag_id=?", (tid,))
    conn.commit()
    conn.close()
    return {"ok": True}


# ---------- 产品 ----------

@app.get("/api/products")
def list_products(
    search: str = "",
    category_id: Optional[int] = None,
    company_id: Optional[int] = None,
    tag_id: Optional[int] = None,
    page: int = 1,
    page_size: int = 20,
    price_field: str = "channel",        # market=市场价 channel=渠道价(进货价)
    price_min: Optional[float] = None,   # 金额区间下限
    price_max: Optional[float] = None,   # 金额区间上限
    sort: str = "",                      # ""=按序号 price_asc/price_desc=按所选价格
):
    conn = get_conn()
    where = []
    args = []
    if search:
        like = f"%{search}%"
        where.append("(p.name LIKE ? OR p.model LIKE ? OR p.intro LIKE ? OR p.params LIKE ? OR p.seq LIKE ?)")
        args += [like, like, like, like, like]
    if category_id:
        where.append("p.category_id=?")
        args.append(category_id)
    if company_id:
        where.append("p.company_id=?")
        args.append(company_id)
    if tag_id:
        where.append("p.id IN (SELECT product_id FROM product_tag WHERE tag_id=?)")
        args.append(tag_id)
    # 金额区间筛选（按所选价格字段）
    pf = "market_price" if price_field == "market" else "channel_price"
    if price_min is not None:
        where.append(f"IFNULL(p.{pf}, -1) >= ?")
        args.append(price_min)
    if price_max is not None:
        where.append(f"IFNULL(p.{pf}, -1) <= ?")
        args.append(price_max)
    # 排序
    if sort in ("price_asc", "price_desc"):
        d = "ASC" if sort == "price_asc" else "DESC"
        order = f"ORDER BY CASE WHEN p.{pf} IS NULL THEN 1 ELSE 0 END, p.{pf} {d}, p.id ASC"
    else:
        order = "ORDER BY CAST(p.seq AS INTEGER) ASC, p.id ASC"

    where_sql = ("WHERE " + " AND ".join(where)) if where else ""
    total = conn.execute(
        f"SELECT COUNT(*) FROM product p {where_sql}", args
    ).fetchone()[0]

    sql = f"""
        SELECT p.*, c.name AS category_name, co.name AS company_name,
               f.filename AS source_filename
        FROM product p
        LEFT JOIN category c ON p.category_id=c.id
        LEFT JOIN company co ON p.company_id=co.id
        LEFT JOIN file f ON p.source_file_id=f.id
        {where_sql}
        {order}
        LIMIT ? OFFSET ?
    """
    rows = conn.execute(sql, args + [page_size, (page - 1) * page_size]).fetchall()
    # 批量取标签和附件（避免 N+1 查询）
    ids = [r["id"] for r in rows]
    tags_map = {}
    atts_map = {}
    if ids:
        ph = ",".join("?" for _ in ids)
        for row in conn.execute(
            f"""SELECT pt.product_id, t.id, t.name FROM product_tag pt
                JOIN tag t ON t.id=pt.tag_id WHERE pt.product_id IN ({ph}) ORDER BY t.name""",
            ids,
        ):
            tags_map.setdefault(row["product_id"], []).append({"id": row["id"], "name": row["name"]})
        for row in conn.execute(
            f"""SELECT id, product_id, filename, kind FROM attachment
                WHERE product_id IN ({ph}) ORDER BY id""",
            ids,
        ):
            kind = "images" if row["kind"] == "image" else "files"
            atts_map.setdefault((row["product_id"], kind), []).append(
                {"id": row["id"], "filename": row["filename"]}
            )
    items = []
    for r in rows:
        d = dict(r)
        d["tags"] = tags_map.get(d["id"], [])
        d["images"] = atts_map.get((d["id"], "images"), [])
        d["files"] = atts_map.get((d["id"], "files"), [])
        items.append(d)
    conn.close()
    return {"total": total, "page": page, "page_size": page_size, "items": items}


@app.get("/api/products/{pid}")
def get_product(pid: int):
    conn = get_conn()
    p = conn.execute(
        "SELECT p.*, c.name AS category_name, co.name AS company_name, "
        "f.filename AS source_filename "
        "FROM product p LEFT JOIN category c ON p.category_id=c.id "
        "LEFT JOIN company co ON p.company_id=co.id "
        "LEFT JOIN file f ON p.source_file_id=f.id WHERE p.id=?", (pid,)
    ).fetchone()
    result = product_full(conn, p)
    conn.close()
    if result is None:
        raise HTTPException(404, "产品不存在")
    return result


class ProductIn(BaseModel):
    seq: Optional[str] = ""
    name: str
    intro: Optional[str] = ""
    params: Optional[str] = ""
    model: Optional[str] = ""
    market_price: Optional[float] = None
    channel_price: Optional[float] = None
    category_name: Optional[str] = ""
    company_name: Optional[str] = ""
    contact_phone: Optional[str] = ""
    contact_person: Optional[str] = ""
    tag_ids: Optional[List[int]] = []


def save_product_tags(conn, product_id, tag_ids):
    conn.execute("DELETE FROM product_tag WHERE product_id=?", (product_id,))
    for t in (tag_ids or []):
        conn.execute("INSERT OR IGNORE INTO product_tag(product_id, tag_id) VALUES(?,?)",
                     (product_id, t))


@app.post("/api/products")
def create_product(p: ProductIn):
    conn = get_conn()
    name = (p.name or "").strip()
    company_name = (p.company_name or "").strip()
    model = (p.model or "").strip()
    # 重复校验：名称 + 公司 + 型号都相同则提示
    dup = conn.execute(
        """SELECT p.id FROM product p JOIN company c ON p.company_id=c.id
           WHERE p.name=? AND c.name=? AND IFNULL(p.model,'')=?""",
        (name, company_name, model),
    ).fetchone()
    if dup:
        conn.close()
        raise HTTPException(409, f"已存在相同产品（名称、公司、型号均相同）：「{name}」-「{company_name}」-「{model or '无型号'}」")
    category_id = get_or_create(conn, "category", p.category_name)
    company_id = get_or_create(conn, "company", p.company_name)
    # 自动生成序号：按录入顺序，全局唯一（取当前最大序号 +1）
    row = conn.execute(
        "SELECT MAX(CAST(seq AS INTEGER)) AS m FROM product WHERE seq GLOB '[0-9]*'"
    ).fetchone()
    seq = str((row["m"] or 0) + 1)
    cur = conn.execute(
        """INSERT INTO product(seq,name,intro,params,model,market_price,channel_price,
           category_id,company_id,contact_phone,contact_person,updated_at)
           VALUES(?,?,?,?,?,?,?,?,?,?,?,?)""",
        (seq, name, p.intro, p.params, model, p.market_price, p.channel_price,
         category_id, company_id, p.contact_phone, p.contact_person, now()),
    )
    pid = cur.lastrowid
    save_product_tags(conn, pid, p.tag_ids)
    conn.commit()
    conn.close()
    return {"id": pid, "seq": seq}


@app.put("/api/products/{pid}")
def update_product(pid: int, p: ProductIn):
    conn = get_conn()
    name = (p.name or "").strip()
    company_name = (p.company_name or "").strip()
    model = (p.model or "").strip()
    # 重复校验（排除自身）
    dup = conn.execute(
        """SELECT p.id FROM product p JOIN company c ON p.company_id=c.id
           WHERE p.name=? AND c.name=? AND IFNULL(p.model,'')=? AND p.id != ?""",
        (name, company_name, model, pid),
    ).fetchone()
    if dup:
        conn.close()
        raise HTTPException(409, f"已存在相同产品（名称、公司、型号均相同）：「{name}」-「{company_name}」-「{model or '无型号'}」")
    category_id = get_or_create(conn, "category", p.category_name)
    company_id = get_or_create(conn, "company", p.company_name)
    # 编辑不改变序号（序号按录入顺序自动生成）
    conn.execute(
        """UPDATE product SET name=?,intro=?,params=?,model=?,market_price=?,
           channel_price=?,category_id=?,company_id=?,contact_phone=?,contact_person=?,updated_at=?
           WHERE id=?""",
        (name, p.intro, p.params, model, p.market_price, p.channel_price,
         category_id, company_id, p.contact_phone, p.contact_person, now(), pid),
    )
    save_product_tags(conn, pid, p.tag_ids)
    conn.commit()
    conn.close()
    return {"ok": True}


@app.delete("/api/products/{pid}")
def delete_product(pid: int):
    conn = get_conn()
    for a in conn.execute("SELECT stored_name, kind FROM attachment WHERE product_id=?", (pid,)).fetchall():
        path = os.path.join(IMAGE_DIR if a["kind"] == "image" else FILE_DIR, a["stored_name"])
        if os.path.exists(path):
            try:
                os.remove(path)
            except OSError:
                pass
    conn.execute("DELETE FROM attachment WHERE product_id=?", (pid,))
    conn.execute("DELETE FROM product_tag WHERE product_id=?", (pid,))
    conn.execute("DELETE FROM product WHERE id=?", (pid,))
    conn.commit()
    conn.close()
    return {"ok": True}


# ---------- 附件 ----------

@app.post("/api/products/{pid}/attachments")
async def upload_attachments(pid: int, files: List[UploadFile] = File(...), kind: str = "auto"):
    conn = get_conn()
    if not conn.execute("SELECT id FROM product WHERE id=?", (pid,)).fetchone():
        conn.close()
        raise HTTPException(404, "产品不存在")
    saved = []
    for f in files:
        ext = os.path.splitext(f.filename or "")[1]
        stored = uuid.uuid4().hex + ext
        ctype = (f.content_type or "").lower()
        fkind = kind if kind in ("image", "file") else ("image" if ctype.startswith("image/") else "file")
        target_dir = IMAGE_DIR if fkind == "image" else FILE_DIR
        path = os.path.join(target_dir, stored)
        data = await f.read()
        with open(path, "wb") as out:
            out.write(data)
        cur = conn.execute(
            "INSERT INTO attachment(product_id, filename, stored_name, kind, size) VALUES(?,?,?,?,?)",
            (pid, f.filename, stored, fkind, len(data)),
        )
        saved.append({"id": cur.lastrowid, "filename": f.filename, "kind": fkind})
    conn.commit()
    conn.close()
    return {"ok": True, "saved": saved}


@app.get("/api/attachments/{aid}/raw")
def raw_attachment(aid: int):
    conn = get_conn()
    a = conn.execute("SELECT * FROM attachment WHERE id=?", (aid,)).fetchone()
    conn.close()
    if a is None:
        raise HTTPException(404, "附件不存在")
    path = os.path.join(IMAGE_DIR if a["kind"] == "image" else FILE_DIR, a["stored_name"])
    if not os.path.exists(path):
        raise HTTPException(404, "文件已丢失")
    return FileResponse(path, filename=a["filename"])


@app.get("/api/attachments/{aid}/download")
def download_attachment(aid: int):
    conn = get_conn()
    a = conn.execute("SELECT * FROM attachment WHERE id=?", (aid,)).fetchone()
    conn.close()
    if a is None:
        raise HTTPException(404, "附件不存在")
    path = os.path.join(IMAGE_DIR if a["kind"] == "image" else FILE_DIR, a["stored_name"])
    if not os.path.exists(path):
        raise HTTPException(404, "文件已丢失")
    return FileResponse(path, filename=a["filename"])


@app.delete("/api/attachments/{aid}")
def delete_attachment(aid: int):
    conn = get_conn()
    a = conn.execute("SELECT * FROM attachment WHERE id=?", (aid,)).fetchone()
    if a:
        path = os.path.join(IMAGE_DIR if a["kind"] == "image" else FILE_DIR, a["stored_name"])
        if os.path.exists(path):
            try:
                os.remove(path)
            except OSError:
                pass
        conn.execute("DELETE FROM attachment WHERE id=?", (aid,))
        conn.commit()
    conn.close()
    return {"ok": True}


# ---------- 文件库 ----------

@app.get("/api/folders")
def list_folders():
    conn = get_conn()
    rows = conn.execute("SELECT * FROM folder ORDER BY id").fetchall()
    conn.close()
    return [dict(r) for r in rows]


@app.post("/api/folders")
def add_folder(payload: dict):
    name = (payload.get("name") or "").strip()
    if not name:
        raise HTTPException(400, "文件夹名不能为空")
    conn = get_conn()
    try:
        cur = conn.execute("INSERT INTO folder(name) VALUES(?)", (name,))
        conn.commit()
        return {"id": cur.lastrowid, "name": name}
    except Exception:
        raise HTTPException(400, "文件夹已存在")
    finally:
        conn.close()


@app.delete("/api/folders/{fid}")
def del_folder(fid: int):
    conn = get_conn()
    row = conn.execute("SELECT name FROM folder WHERE id=?", (fid,)).fetchone()
    if row:
        files = conn.execute("SELECT * FROM file WHERE folder=?", (row["name"],)).fetchall()
        for f in files:
            path = os.path.join(LIBRARY_DIR, f["stored_name"])
            if os.path.exists(path):
                try:
                    os.remove(path)
                except OSError:
                    pass
        conn.execute("DELETE FROM file WHERE folder=?", (row["name"],))
        conn.execute("DELETE FROM folder WHERE id=?", (fid,))
        conn.commit()
    conn.close()
    return {"ok": True}


@app.get("/api/files")
def list_files(search: str = "", folder: Optional[str] = None):
    conn = get_conn()
    where = []
    args = []
    if search:
        where.append("filename LIKE ?")
        args.append(f"%{search}%")
    if folder is not None and folder != "__all__":
        where.append("folder = ?")
        args.append(folder)
    where_sql = ("WHERE " + " AND ".join(where)) if where else ""
    rows = conn.execute(f"SELECT * FROM file {where_sql} ORDER BY id DESC", args).fetchall()
    conn.close()
    return [dict(r) for r in rows]


@app.post("/api/files")
async def upload_file(file: UploadFile = File(...), folder: str = Form("")):
    conn = get_conn()
    existing = conn.execute(
        "SELECT id FROM file WHERE filename=? AND folder=?", (file.filename, folder)
    ).fetchone()
    if existing:
        conn.close()
        raise HTTPException(409, f"文件「{file.filename}」已存在，已跳过")
    ext = os.path.splitext(file.filename or "")[1]
    stored = uuid.uuid4().hex + ext
    path = os.path.join(LIBRARY_DIR, stored)
    data = await file.read()
    with open(path, "wb") as out:
        out.write(data)
    cur = conn.execute(
        "INSERT INTO file(filename, stored_name, folder, size) VALUES(?,?,?,?)",
        (file.filename, stored, folder, len(data)),
    )
    conn.commit()
    conn.close()
    return {"id": cur.lastrowid, "filename": file.filename, "folder": folder, "size": len(data)}


@app.get("/api/files/{fid}/download")
def download_file(fid: int):
    conn = get_conn()
    f = conn.execute("SELECT * FROM file WHERE id=?", (fid,)).fetchone()
    conn.close()
    if f is None:
        raise HTTPException(404, "文件不存在")
    path = os.path.join(LIBRARY_DIR, f["stored_name"])
    if not os.path.exists(path):
        raise HTTPException(404, "文件已丢失")
    return FileResponse(path, filename=f["filename"])


@app.delete("/api/files/{fid}")
def delete_file(fid: int):
    conn = get_conn()
    f = conn.execute("SELECT * FROM file WHERE id=?", (fid,)).fetchone()
    if f:
        path = os.path.join(LIBRARY_DIR, f["stored_name"])
        if os.path.exists(path):
            try:
                os.remove(path)
            except OSError:
                pass
        conn.execute("UPDATE product SET source_file_id=NULL WHERE source_file_id=?", (fid,))
        conn.execute("DELETE FROM file WHERE id=?", (fid,))
        conn.commit()
    conn.close()
    return {"ok": True}


# ---------- 文件库在线预览 ----------

from html import escape as _esc
import csv as _csv
import io as _io

# 浏览器原生可预览(图片 / 音视频 / PDF):直接以 inline 形式返回
_PREVIEW_NATIVE = {
    ".pdf": "application/pdf",
    ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
    ".gif": "image/gif", ".webp": "image/webp", ".bmp": "image/bmp",
    ".svg": "image/svg+xml", ".ico": "image/x-icon",
    ".mp4": "video/mp4", ".webm": "video/webm", ".mov": "video/quicktime",
    ".mp3": "audio/mpeg", ".wav": "audio/wav", ".ogg": "audio/ogg",
    ".m4a": "audio/mp4", ".aac": "audio/aac", ".flac": "audio/flac",
}
# 纯文本直显
_PREVIEW_TEXT = {
    ".txt", ".md", ".markdown", ".log", ".json", ".xml",
    ".ini", ".cfg", ".yaml", ".yml", ".py", ".js", ".css", ".html", ".htm",
}
# 表格(Excel / CSV)
_PREVIEW_TABLE = {".xlsx", ".xlsm", ".xls", ".csv"}
# Word / PPT
_PREVIEW_DOCX = {".docx"}
_PREVIEW_PPTX = {".pptx"}
# Office 原生文档:可调用本机 Office/WPS COM 保真转 PDF
_PREVIEW_OFFICE = {".docx", ".doc", ".xlsx", ".xlsm", ".xls", ".pptx", ".ppt"}


def _read_text_smart(path, limit=8 * 1024 * 1024):
    with open(path, "rb") as fh:
        raw = fh.read(limit)
    for enc in ("utf-8", "gb18030", "utf-16", "latin-1"):
        try:
            return raw.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return raw.decode("utf-8", errors="replace")


def _page_head(title):
    return ("<!DOCTYPE html><html lang='zh-CN'><head><meta charset='utf-8'>"
            "<title>文件预览</title><style>"
            "body{font-family:'Microsoft YaHei',system-ui,sans-serif;margin:0;background:#fff;color:#1f2937}"
            ".top{position:sticky;top:0;background:#f9fafb;border-bottom:1px solid #e5e7eb;padding:10px 18px;font-size:14px;color:#374151;display:flex;align-items:center;gap:8px}"
            ".top b{max-width:60vw;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;display:inline-block;vertical-align:bottom}"
            ".body{padding:18px;overflow:auto}"
            "pre.txt{white-space:pre-wrap;word-break:break-all;font-family:Consolas,'Courier New',monospace;font-size:13px;line-height:1.7;margin:0}"
            "table{border-collapse:collapse;width:100%;font-size:13px;margin-bottom:20px;background:#fff}"
            "th,td{border:1px solid #d1d5db;padding:6px 10px;text-align:left;word-break:break-all}"
            "th{background:#f3f4f6;font-weight:600;position:sticky;top:41px}"
            "td.num,th.num{text-align:right}"
            ".sheet-tab{display:inline-block;padding:6px 14px;margin:4px 6px 4px 0;border:1px solid #d1d5db;border-radius:16px;cursor:pointer;font-size:13px;color:#374151;background:#fff}"
            ".sheet-tab.on{background:#2563eb;color:#fff;border-color:#2563eb}"
            ".sheet-panel{display:none}.sheet-panel.on{display:block}"
            ".hint{color:#6b7280;font-size:13px;padding:30px;text-align:center}"
            ".p,.doc-p{white-space:pre-wrap;word-break:break-word;margin:0 0 10px;line-height:1.75;font-size:14px}"
            "h1.dl{font-size:17px;margin:4px 0 14px;color:#111827}"
            ".slide-card{border:1px solid #e5e7eb;border-radius:8px;padding:14px 16px;margin-bottom:14px;background:#fcfcfd}"
            ".slide-no{color:#6b7280;font-size:12px;margin-bottom:8px}"
            "</style></head><body>"
            f"<div class='top'>📄 在线预览 · <b>{_esc(title)}</b>"
            "&nbsp;<a style='color:#2563eb;text-decoration:none;margin-left:auto' href='#' onclick='window.close();return false'>关闭窗口</a></div>")


def _page_end():
    return "</body></html>"


# ---------- 异步预览任务框架 ----------

PREVIEW_DIR = os.path.join(BASE_DIR, "uploads", "preview")
os.makedirs(PREVIEW_DIR, exist_ok=True)
PREVIEW_THRESHOLD = 30 * 1024 * 1024  # 30MB,超过则后台异步生成

_preview_tasks = {}  # task_id -> dict(status, fid, filename, started, finished, error, size, result_path)


def _new_task(fid, filename):
    tid = uuid.uuid4().hex
    _preview_tasks[tid] = {
        "status": "queued",
        "fid": fid,
        "filename": filename,
        "started": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "finished": None,
        "error": None,
        "size": 0,
        "result_path": None,
    }
    return tid


def _task_view(tid):
    d = _preview_tasks.get(tid)
    if not d:
        return {"status": "not_found", "error": "任务不存在或已过期"}
    v = {"status": d["status"], "filename": d["filename"], "error": d["error"],
         "note": d.get("note")}
    if d["status"] == "done":
        v["url"] = f"/api/preview/result/{tid}"
        v["kind"] = "pdf" if (d.get("result_path") or "").lower().endswith(".pdf") else "html"
    return v


def _render_friendly_html(name, message, fid=None, status=None):
    body = (f"<div class='hint'>⚠️ {_esc(message)}<br>"
           f"<div style='margin-top:14px;color:#6b7280;font-size:13px'>文件: <b>{_esc(name or '')}</b></div>")
    if fid is not None:
        body += (f"<div style='margin-top:18px'><a href='/api/files/{fid}/download' "
                 f"style='color:#2563eb;text-decoration:none;font-size:14px'>⬇ 下载原文件</a></div>")
    body += "</div>"
    return _page_head(f"无法预览 · {name or ''}") + body + _page_end()


def _render_preview_html(path, ext, filename):
    """按扩展名分发到对应 HTML 渲染器(集成原表格/docx/pptx/文本)。"""
    if ext in _PREVIEW_TABLE:
        return _build_table_html(path, ext, filename)
    if ext in _PREVIEW_DOCX:
        return _build_docx_html(path, filename)
    if ext in _PREVIEW_PPTX:
        return _build_pptx_html(path, filename)
    if ext in _PREVIEW_TEXT:
        text = _read_text_smart(path)
        body = f"<pre class='txt'>{_esc(text)}</pre>"
        return _page_head(filename) + "<div class='body'>" + body + "</div>" + _page_end()
    return None


OFFICE_CONVERT_TIMEOUT = 180  # Office 转换最长等待(秒),超时视为失败并终止进程


def _convert_office_to_pdf(src, dst, ext):
    """以子进程方式调用本机 Office / WPS COM 保真转 PDF。

    好处:COM 卡死/超时只影响子进程,可被主服务强制终止,不会拖死预览服务线程。
    """
    cli = os.path.join(BASE_DIR, "office_convert_cli.py")
    py = sys.executable
    import subprocess as _sp
    proc = None
    try:
        proc = _sp.Popen(
            [py, cli, src, dst, ext],
            stdout=_sp.PIPE, stderr=_sp.STDOUT, text=True, encoding="utf-8", errors="replace",
        )
        try:
            out, _ = proc.communicate(timeout=OFFICE_CONVERT_TIMEOUT)
        except _sp.TimeoutExpired:
            proc.kill()
            try:
                proc.communicate(timeout=5)
            except Exception:
                pass
            raise RuntimeError(f"Office 转换超时(>{OFFICE_CONVERT_TIMEOUT}s),文件过大或电脑性能不足")
        if proc.returncode != 0 or not (os.path.exists(dst) and os.path.getsize(dst) > 1000):
            raise RuntimeError((out or "").strip()[-200:] or f"Office 转换失败(exit={proc.returncode})")
    finally:
        if proc is not None and proc.poll() is None:
            try:
                proc.kill()
            except Exception:
                pass


def _run_preview_task(tid, path, ext, name):
    t = _preview_tasks[tid]
    t["status"] = "running"
    t["converting"] = ext in _PREVIEW_OFFICE
    out_stem = os.path.join(PREVIEW_DIR, tid)
    try:
        if ext in _PREVIEW_OFFICE:
            # 保真路线:本机 Office/WPS 转 PDF
            try:
                _convert_office_to_pdf(path, out_stem + ".pdf", ext)
                t["result_path"] = out_stem + ".pdf"
                t["status"] = "done"
                t["converted_as"] = "pdf"
                t["note"] = None
            except Exception as conv_e:
                # COM 转换失败(未装 Office / 文件损坏 / 超大超时等)→ 尝试降级为提取式预览
                try:
                    html = _render_preview_html(path, ext, name)
                    if html is None:
                        raise RuntimeError("fallback_none")
                    with open(out_stem + ".html", "w", encoding="utf-8") as f:
                        f.write(html)
                    t["result_path"] = out_stem + ".html"
                    t["status"] = "done"
                    t["converted_as"] = "html"
                    t["note"] = f"本机保真转换未成功(原因:{str(conv_e)[:80]}),已切换为文本预览"
                except Exception:
                    # 连提取式也不支持该类型时,如实报告转换原因,避免误导为"不支持预览"
                    t["status"] = "failed"
                    t["error"] = f"保真转换失败:{str(conv_e)[:120]}。请下载原文件,用本机 Office/WPS 打开查看"
        else:
            html = _render_preview_html(path, ext, name)
            if html is None:
                raise RuntimeError("暂不支持此类型的在线预览")
            with open(out_stem + ".html", "w", encoding="utf-8") as f:
                f.write(html)
            t["result_path"] = out_stem + ".html"
            t["status"] = "done"
    except Exception as e:
        t["status"] = "failed"
        t["error"] = str(e)
    finally:
        t["finished"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")


@app.get("/api/files/{fid}/preview")
def preview_file(fid: int):
    """兼容旧接口:文件不存在/类型不支持时返回友好 HTML(不再让 iframe 出 localhost 拒绝连接)。"""
    conn = get_conn()
    f = conn.execute("SELECT * FROM file WHERE id=?", (fid,)).fetchone()
    conn.close()
    name = f["filename"] if f else ""
    if f is None:
        return Response(content=_render_friendly_html(name, "该文件不存在或已被删除", fid=None),
                        media_type="text/html; charset=utf-8", status_code=404)
    path = os.path.join(LIBRARY_DIR, f["stored_name"])
    if not os.path.exists(path):
        return Response(content=_render_friendly_html(name, "文件已丢失,请重新上传", fid=fid),
                        media_type="text/html; charset=utf-8", status_code=404)
    ext = os.path.splitext(f["filename"])[1].lower()
    if ext in _PREVIEW_NATIVE:
        return FileResponse(
            path, media_type=_PREVIEW_NATIVE[ext],
            headers={"Content-Disposition": f"inline; filename*=UTF-8''{quote(f['filename'])}"},
        )
    try:
        html = _render_preview_html(path, ext, f["filename"])
        if html is None:
            return Response(content=_render_friendly_html(f["filename"], "暂不支持此文件类型的在线预览,请下载后查看", fid=fid),
                            media_type="text/html; charset=utf-8", status_code=415)
        return Response(content=html, media_type="text/html; charset=utf-8")
    except Exception as e:
        return Response(content=_render_friendly_html(f["filename"], f"解析失败: {e}", fid=fid),
                        media_type="text/html; charset=utf-8", status_code=500)


# 原生类型 inline 预览(给 PDF / 图片 / 音视频 用;浏览器原生渲染)
@app.get("/api/files/{fid}/inline")
def preview_inline(fid: int):
    conn = get_conn()
    f = conn.execute("SELECT * FROM file WHERE id=?", (fid,)).fetchone()
    conn.close()
    if f is None:
        raise HTTPException(404, "文件不存在")
    path = os.path.join(LIBRARY_DIR, f["stored_name"])
    if not os.path.exists(path):
        raise HTTPException(404, "文件已丢失")
    ext = os.path.splitext(f["filename"])[1].lower()
    if ext not in _PREVIEW_NATIVE:
        raise HTTPException(400, "此类型不支持 inline 预览")
    return FileResponse(
        path, media_type=_PREVIEW_NATIVE[ext],
        headers={"Content-Disposition": f"inline; filename*=UTF-8''{quote(f['filename'])}"},
    )


# 异步预览任务(大文件后台生成,前端轮询状态)
@app.post("/api/files/{fid}/preview/async")
def preview_async(fid: int, bg: BackgroundTasks):
    conn = get_conn()
    f = conn.execute("SELECT * FROM file WHERE id=?", (fid,)).fetchone()
    conn.close()
    if f is None:
        raise HTTPException(404, "文件不存在")
    ext = os.path.splitext(f["filename"])[1].lower()
    size = os.path.getsize(os.path.join(LIBRARY_DIR, f["stored_name"])) if os.path.exists(os.path.join(LIBRARY_DIR, f["stored_name"])) else 0
    # 原生类型(PDF/图片/音视频)浏览器直接展示,无需异步
    if ext in _PREVIEW_NATIVE:
        if not os.path.exists(os.path.join(LIBRARY_DIR, f["stored_name"])):
            raise HTTPException(404, "文件已丢失")
        return {"task_id": "native", "status": "done",
                "url": f"/api/files/{fid}/inline",
                "filename": f["filename"], "size": size}
    if not os.path.exists(os.path.join(LIBRARY_DIR, f["stored_name"])):
        raise HTTPException(404, "文件已丢失")
    tid = _new_task(fid, f["filename"])
    _preview_tasks[tid]["size"] = size
    full_path = os.path.join(LIBRARY_DIR, f["stored_name"])
    if ext in _PREVIEW_OFFICE:
        # Office 保真转换(COM 需数秒~十几秒):一律后台执行,避免请求阻塞
        bg.add_task(_run_preview_task, tid, full_path, ext, f["filename"])
    elif size <= PREVIEW_THRESHOLD:
        # 文本/表格类小文件:直接同步生成(几十毫秒即可完成)
        _run_preview_task(tid, full_path, ext, f["filename"])
    else:
        bg.add_task(_run_preview_task, tid, full_path, ext, f["filename"])
    v = _task_view(tid)
    return {"task_id": tid, "status": v["status"],
            "filename": f["filename"], "size": size,
            "url": v.get("url"), "note": v.get("note")}


@app.get("/api/preview/tasks/{tid}")
def preview_task_status(tid: str):
    return _task_view(tid)


@app.get("/api/preview/result/{tid}")
def preview_task_result(tid: str):
    d = _preview_tasks.get(tid)
    if not d or d["status"] != "done" or not d["result_path"]:
        raise HTTPException(404, "预览尚未就绪")
    rp = d["result_path"]
    if rp.lower().endswith(".pdf"):
        return FileResponse(
            rp, media_type="application/pdf",
            headers={"Content-Disposition": f"inline; filename*=UTF-8''{quote(d['filename'] + '.pdf')}"},
        )
    return FileResponse(rp, media_type="text/html; charset=utf-8")


def _fmt_cell(v):
    if v is None:
        return ""
    if isinstance(v, float):
        return ("%g" % v)
    return str(v)


def _rows_to_table_html(rows, sticky_header=True):
    if not rows:
        return "<p class='hint'>空表格</p>"
    head = "".join(f"<th>{_esc(str(c))}</th>" for c in rows[0]) if rows[0] else ""
    trs = []
    for r in rows[1:]:
        tds = "".join(f"<td>{_esc(_fmt_cell(c))}</td>" for c in r)
        trs.append(f"<tr>{tds}</tr>")
    return f"<table><thead><tr>{head}</tr></thead><tbody>{''.join(trs)}</tbody></table>"


def _build_table_html(path, ext, filename):
    title = f"表格式文件 · {filename}"
    try:
        if ext == ".csv":
            text = _read_text_smart(path)
            try:
                reader = list(_csv.reader(_io.StringIO(text)))
            except Exception:
                reader = [line.split(",") for line in text.splitlines() if line.strip()]
            sheets = [("数据", reader)]
        else:
            from openpyxl import load_workbook
            wb = load_workbook(path, data_only=True, read_only=True)
            sheets = [(ws.title, [list(r) for r in ws.iter_rows(values_only=True)]) for ws in wb.worksheets]
    except Exception as e:
        return (_page_head(title)
                + f"<div class='hint'>解析失败：{_esc(str(e))}（建议直接下载原文件）</div>" + _page_end())
    if len(sheets) == 1:
        tabs = ""
        panels = "<div class='body'>" + _rows_to_table_html(sheets[0][1]) + "</div>"
    else:
        tabs = "".join(f"<span class='sheet-tab' data-i='{i}'>{_esc(name)}</span>" for i, (name, _) in enumerate(sheets))
        panels = "".join(
            f"<div class='sheet-panel' data-p='{i}'>" + _rows_to_table_html(rows) + "</div>"
            for i, (_, rows) in enumerate(sheets))
        tabs = ("<div style='padding:8px 16px 0;position:sticky;top:41px;background:#fff;z-index:2;border-bottom:1px solid #eee'>"
                + tabs + "</div>")
        tabs += ("<script>"
                 "(function(){var a=document.querySelectorAll('.sheet-tab'),p=document.querySelectorAll('.sheet-panel');"
                 "function on(i){a.forEach(function(x){x.classList.toggle('on',+x.dataset.i===i)});"
                 "p.forEach(function(x){x.classList.toggle('on',+x.dataset.p===i)})}on(0);"
                 "a.forEach(function(x){x.onclick=function(){on(+x.dataset.i)}})})();"
                 "</script>")
    return _page_head(title) + tabs + panels + _page_end()


def _build_docx_html(path, filename):
    title = f"Word 文档 · {filename}"
    try:
        import docx
        from docx.table import Table as _DocxTable
        from docx.text.paragraph import Paragraph as _DocxPara
        from docx.oxml.ns import qn

        d = docx.Document(path)
        parts = []

        def style_of(p):
            try:
                return (p.style.name or "") if p.style else ""
            except Exception:
                return ""

        for child in d.element.body.iterchildren():
            if child.tag == qn("w:p"):
                p = _DocxPara(child, d)
                txt = p.text.strip()
                if txt:
                    st = style_of(p)
                    if "Heading" in st or "标题" in st:
                        lvl = "".join(ch for ch in st if ch.isdigit()) or "1"
                        parts.append(f"<h{min(int(lvl) + 1, 6) if lvl else 3} class='dl' style='margin-top:18px'>{_esc(txt)}</h{min(int(lvl) + 1, 6) if lvl else 3}>")
                    else:
                        parts.append(f"<p class='doc-p'>{_esc(txt)}</p>")
            elif child.tag == qn("w:tbl"):
                t = _DocxTable(child, d)
                rows = [[cell.text.strip() for cell in r.cells] for r in t.rows]
                if rows:
                    parts.append(_rows_to_table_html(rows))
        body = "".join(parts) if parts else "<div class='hint'>该文档没有可预览的正文内容</div>"
    except Exception as e:
        return (_page_head(title) + f"<div class='hint'>解析失败：{_esc(str(e))}（建议直接下载原文件查看）</div>" + _page_end())
    return _page_head(title) + "<div class='body'>" + body + "</div>" + _page_end()


def _build_pptx_html(path, filename):
    title = f"PPT 演示文稿 · {filename}"
    try:
        from pptx import Presentation
        prs = Presentation(path)
        cards = []
        for idx, slide in enumerate(prs.slides, start=1):
            blocks = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    txt = "\n".join(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip())
                    if txt:
                        blocks.append(f"<p class='doc-p'>{_esc(txt)}</p>")
                elif getattr(shape, "has_table", False):
                    tbl = shape.table
                    rows = [[c.text.strip() for c in row.cells] for row in tbl.rows]
                    if rows:
                        blocks.append(_rows_to_table_html(rows))
            if blocks:
                cards.append(f"<div class='slide-card'><div class='slide-no'>第 {idx} 页</div>{''.join(blocks)}</div>")
        body = "".join(cards) if cards else "<div class='hint'>该演示文稿没有可预览的文字内容</div>"
    except Exception as e:
        return (_page_head(title) + f"<div class='hint'>解析失败：{_esc(str(e))}（建议直接下载原文件查看）</div>" + _page_end())
    return _page_head(title) + "<div class='body'>" + body + "</div>" + _page_end()
def batch_download(payload: dict):
    ids = payload.get("ids") or []
    if not ids:
        raise HTTPException(400, "未选择文件")
    conn = get_conn()
    placeholders = ",".join("?" for _ in ids)
    files = conn.execute(f"SELECT * FROM file WHERE id IN ({placeholders})", ids).fetchall()
    conn.close()
    if not files:
        raise HTTPException(404, "文件不存在")
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for f in files:
            path = os.path.join(LIBRARY_DIR, f["stored_name"])
            if os.path.exists(path):
                arcname = f["filename"]
                if f["folder"]:
                    arcname = f"{f['folder']}/{f['filename']}"
                zf.write(path, arcname)
    fname = f"文件下载_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip"
    headers = {"Content-Disposition": f"attachment; filename*=UTF-8''{quote(fname)}"}
    return Response(content=zip_buf.getvalue(), media_type="application/zip", headers=headers)


@app.post("/api/files/batch-move")
def batch_move(payload: dict):
    ids = payload.get("ids") or []
    folder = payload.get("folder") or ""
    if not ids:
        raise HTTPException(400, "未选择文件")
    conn = get_conn()
    placeholders = ",".join("?" for _ in ids)
    conn.execute(f"UPDATE file SET folder=? WHERE id IN ({placeholders})", [folder] + ids)
    conn.commit()
    conn.close()
    return {"ok": True, "moved": len(ids)}


# ---------- 导出 / 导入 ----------

HEADERS = ["序号", "名称", "介绍", "参数", "型号", "市场价", "渠道价",
           "产品类型", "公司名称", "联系方式", "联系人", "标签"]


@app.get("/api/export")
def export_products(
    search: str = "",
    category_id: Optional[int] = None,
    company_id: Optional[int] = None,
    tag_id: Optional[int] = None,
):
    conn = get_conn()
    where = []
    args = []
    if search:
        like = f"%{search}%"
        where.append("(p.name LIKE ? OR p.model LIKE ? OR p.intro LIKE ? OR p.params LIKE ?)")
        args += [like, like, like, like]
    if category_id:
        where.append("p.category_id=?"); args.append(category_id)
    if company_id:
        where.append("p.company_id=?"); args.append(company_id)
    if tag_id:
        where.append("p.id IN (SELECT product_id FROM product_tag WHERE tag_id=?)"); args.append(tag_id)
    where_sql = ("WHERE " + " AND ".join(where)) if where else ""
    rows = conn.execute(
        f"""SELECT p.*, c.name AS category_name, co.name AS company_name
            FROM product p LEFT JOIN category c ON p.category_id=c.id
            LEFT JOIN company co ON p.company_id=co.id {where_sql}
            ORDER BY CAST(p.seq AS INTEGER) ASC, p.id ASC""",
        args,
    ).fetchall()

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "产品"
    ws.append(HEADERS)
    for c in ws[1]:
        c.font = Font(bold=True)
    for r in rows:
        tags = get_product_tags(conn, r["id"])
        tag_text = "、".join(t["name"] for t in tags)
        ws.append([
            r["seq"], r["name"], r["intro"], r["params"], r["model"],
            r["market_price"], r["channel_price"], r["category_name"],
            r["company_name"], r["contact_phone"], r["contact_person"], tag_text,
        ])
    conn.close()

    buf = io.BytesIO()
    wb.save(buf)
    fname = f"产品导出_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
    headers = {"Content-Disposition": f"attachment; filename*=UTF-8''{quote(fname)}"}
    return Response(
        content=buf.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers=headers,
    )


@app.get("/api/export/template")
def export_template():
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "产品导入模板"
    ws.append(HEADERS)
    for c in ws[1]:
        c.font = Font(bold=True)
    ws.append(["001", "示例产品", "产品介绍", "参数说明", "A-100", 100, 80,
               "电子产品", "示例公司", "13800000000", "张三", "新品、热销"])
    buf = io.BytesIO()
    wb.save(buf)
    headers = {"Content-Disposition": f"attachment; filename*=UTF-8''{quote('product_template.xlsx')}"}
    return Response(
        content=buf.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers=headers,
    )


def get_or_create(conn, table, name):
    name = (name or "").strip()
    if not name:
        return None
    row = conn.execute(f"SELECT id FROM {table} WHERE name=?", (name,)).fetchone()
    if row:
        return row["id"]
    cur = conn.execute(f"INSERT INTO {table}(name) VALUES(?)", (name,))
    return cur.lastrowid


def safe_sheet_name(name):
    for ch in ["\\", "/", "?", "*", "[", "]", ":"]:
        name = name.replace(ch, "")
    return (name or "未分类")[:31]


@app.post("/api/export/selected")
def export_selected(payload: dict):
    ids = payload.get("ids") or []
    if not ids:
        raise HTTPException(400, "未选择任何产品")
    conn = get_conn()
    placeholders = ",".join("?" for _ in ids)
    rows = conn.execute(
        f"""SELECT p.*, c.name AS category_name, co.name AS company_name
            FROM product p LEFT JOIN category c ON p.category_id=c.id
            LEFT JOIN company co ON p.company_id=co.id
            WHERE p.id IN ({placeholders}) ORDER BY c.name, CAST(p.seq AS INTEGER), p.id""",
        ids,
    ).fetchall()
    prod_info = {}
    for r in rows:
        pid = r["id"]
        tags = [t["name"] for t in get_product_tags(conn, pid)]
        atts = conn.execute("SELECT * FROM attachment WHERE product_id=?", (pid,)).fetchall()
        prod_info[pid] = {
            "tags": tags,
            "images": [a for a in atts if a["kind"] == "image"],
            "files": [a for a in atts if a["kind"] == "file"],
        }
    conn.close()

    # 前端传来的购物车价格覆盖（id -> {market_price, channel_price}），不落库
    overrides = {int(k): v for k, v in (payload.get("price_overrides") or {}).items()}

    groups = {}
    for r in rows:
        cat = r["category_name"] or "未分类"
        groups.setdefault(cat, []).append(r)

    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    img_col = 13  # 图片列（第13列 = M）
    img_letter = openpyxl.utils.get_column_letter(img_col)
    grand_market = 0.0
    grand_channel = 0.0
    grand_count = 0
    for cat, items in groups.items():
        ws = wb.create_sheet(title=safe_sheet_name(cat))
        ws.append(["序号", "名称", "型号", "市场价", "渠道价", "产品类型",
                   "公司", "联系人", "联系方式", "标签", "介绍", "参数", "图片"])
        for c in ws[1]:
            c.font = Font(bold=True)
        cat_market = 0.0
        cat_channel = 0.0
        for idx, r in enumerate(items, start=2):
            info = prod_info[r["id"]]
            ov = overrides.get(r["id"], {})
            mp = ov.get("market_price", r["market_price"])
            cp = ov.get("channel_price", r["channel_price"])
            try: mp = float(mp) if mp is not None else None
            except (TypeError, ValueError): mp = r["market_price"]
            try: cp = float(cp) if cp is not None else None
            except (TypeError, ValueError): cp = r["channel_price"]
            if mp is not None: cat_market += mp
            if cp is not None: cat_channel += cp
            ws.append([
                r["seq"], r["name"], r["model"], mp, cp,
                r["category_name"], r["company_name"], r["contact_person"], r["contact_phone"],
                "、".join(info["tags"]), r["intro"], r["params"],
            ])
            imgs = info["images"]
            if imgs:
                path = os.path.join(IMAGE_DIR, imgs[0]["stored_name"])
                if os.path.exists(path):
                    try:
                        ximg = XLImage(path)
                        ximg.width = 90
                        ximg.height = 90
                        ws.add_image(ximg, f"{img_letter}{idx}")
                    except Exception:
                        pass
            ws.row_dimensions[idx].height = 70
        # 分组合计行
        sum_row = ws.max_row + 1
        ws.cell(row=sum_row, column=2, value="小计").font = Font(bold=True)
        ws.cell(row=sum_row, column=4, value=round(cat_market, 2)).font = Font(bold=True)
        ws.cell(row=sum_row, column=5, value=round(cat_channel, 2)).font = Font(bold=True)
        ws.cell(row=sum_row, column=6, value=f"{len(items)} 个产品").font = Font(bold=True)
        grand_market += cat_market
        grand_channel += cat_channel
        grand_count += len(items)
        widths = [8, 28, 16, 12, 12, 16, 18, 10, 14, 16, 40, 40, 14]
        for i, w in enumerate(widths, 1):
            ws.column_dimensions[openpyxl.utils.get_column_letter(i)].width = w

    # 总计 sheet
    ws_total = wb.create_sheet(title="总计")
    ws_total.append(["项目", "数值"])
    ws_total.append(["产品总数（个）", grand_count])
    ws_total.append(["市场价合计（元）", round(grand_market, 2)])
    ws_total.append(["渠道价合计（元）", round(grand_channel, 2)])
    for c in ws_total[1]:
        c.font = Font(bold=True)
    for r in range(2, 5):
        ws_total.cell(row=r, column=1).font = Font(bold=True)
    ws_total.column_dimensions["A"].width = 22
    ws_total.column_dimensions["B"].width = 18

    excel_buf = io.BytesIO()
    wb.save(excel_buf)

    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("产品导出.xlsx", excel_buf.getvalue())
        for r in rows:
            files = prod_info[r["id"]]["files"]
            if not files:
                continue
            name_part = re.sub(r'[\\/:*?"<>|]', "", r["name"])[:30]
            for f in files:
                path = os.path.join(FILE_DIR, f["stored_name"])
                if os.path.exists(path):
                    arcname = f"附件/{r['seq']}_{name_part}_{f['filename']}"
                    zf.write(path, arcname)

    fname = f"选购产品_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip"
    headers = {"Content-Disposition": f"attachment; filename*=UTF-8''{quote(fname)}"}
    return Response(
        content=zip_buf.getvalue(),
        media_type="application/zip",
        headers=headers,
    )


@app.post("/api/import")
async def import_products(file: UploadFile = File(...)):
    data = await file.read()
    try:
        wb = openpyxl.load_workbook(io.BytesIO(data))
    except Exception:
        raise HTTPException(400, "无法解析文件，请上传 .xlsx 格式")
    ws = wb.active
    rows = list(ws.iter_rows(values_only=True))
    if not rows:
        raise HTTPException(400, "文件为空")
    header = [str(h).strip() if h else "" for h in rows[0]]
    idx = {h: i for i, h in enumerate(header)}

    def col(name):
        i = idx.get(name)
        return rows[r][i] if i is not None and i < len(rows[r]) else None

    conn = get_conn()
    created = 0
    errors = []
    for r in range(1, len(rows)):
        name = (col("名称") or "")
        if name is None or str(name).strip() == "":
            continue
        name = str(name).strip()
        category_id = get_or_create(conn, "category", str(col("产品类型") or ""))
        company_id = get_or_create(conn, "company", str(col("公司名称") or ""))
        try:
            cur = conn.execute(
                """INSERT INTO product(seq,name,intro,params,model,market_price,channel_price,
                   category_id,company_id,contact_phone,contact_person,updated_at)
                   VALUES(?,?,?,?,?,?,?,?,?,?,?,?)""",
                (str(col("序号") or ""), name, str(col("介绍") or ""), str(col("参数") or ""),
                 str(col("型号") or ""), col("市场价"), col("渠道价"), category_id, company_id,
                 str(col("联系方式") or ""), str(col("联系人") or ""), now()),
            )
            pid = cur.lastrowid
            tag_text = str(col("标签") or "")
            for t in [x for x in re_split_tags(tag_text) if x]:
                tid = get_or_create(conn, "tag", t)
                if tid:
                    conn.execute("INSERT OR IGNORE INTO product_tag(product_id, tag_id) VALUES(?,?)", (pid, tid))
            created += 1
        except Exception as e:
            errors.append(f"第{r+1}行: {e}")
    conn.commit()
    conn.close()
    return {"created": created, "errors": errors}


def re_split_tags(text):
    import re
    return re.split(r"[、,，;；/\\|]", text)


# ---------- 统计 ----------

@app.get("/api/stats")
def get_stats():
    conn = get_conn()
    total = conn.execute("SELECT COUNT(*) FROM product").fetchone()[0]
    by_category = [dict(r) for r in conn.execute(
        "SELECT c.name AS name, COUNT(p.id) AS value FROM category c "
        "LEFT JOIN product p ON p.category_id=c.id GROUP BY c.id ORDER BY value DESC").fetchall()]
    by_company = [dict(r) for r in conn.execute(
        "SELECT co.name AS name, COUNT(p.id) AS value FROM company co "
        "LEFT JOIN product p ON p.company_id=co.id GROUP BY co.id ORDER BY value DESC").fetchall()]
    by_tag = [dict(r) for r in conn.execute(
        "SELECT t.name AS name, COUNT(pt.product_id) AS value FROM tag t "
        "LEFT JOIN product_tag pt ON pt.tag_id=t.id GROUP BY t.id ORDER BY value DESC").fetchall()]
    conn.close()
    return {"total": total, "by_category": by_category, "by_company": by_company, "by_tag": by_tag}


# ---------- 备份 ----------

@app.get("/api/backup")
def backup():
    if not os.path.exists(DB_PATH):
        raise HTTPException(404, "数据库文件不存在")
    fname = f"产品库备份_{datetime.now().strftime('%Y%m%d_%H%M%S')}.db"
    return FileResponse(DB_PATH, filename=fname, media_type="application/octet-stream")


# ---------- 前端静态 ----------
# 关闭浏览器缓存，确保改代码后用户刷新即可拿到最新版本
@app.middleware("http")
async def no_cache_static(request, call_next):
    response = await call_next(request)
    if not request.url.path.startswith("/api"):
        response.headers["Cache-Control"] = "no-cache, no-store, must-revalidate"
        response.headers["Pragma"] = "no-cache"
        response.headers["Expires"] = "0"
    return response


app.mount("/", StaticFiles(directory=os.path.join(BASE_DIR, "static"), html=True), name="static")


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
